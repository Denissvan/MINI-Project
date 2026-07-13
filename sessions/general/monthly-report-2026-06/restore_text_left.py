from copy import deepcopy
from pathlib import Path
import shutil
from pptx import Presentation

current = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹.pptx')
source = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹_backup_20260623_204031.pptx')
backup = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹.backup-before-restore-text-left.pptx')
shutil.copy2(current, backup)

cur = Presentation(str(current))
src = Presentation(str(source))
cur_slide = cur.slides[1]
src_slide = src.slides[1]

keep_names = {'文本框 1', '直接连接符 62', '矩形 66'}
# Delete current left-side redesigned shapes. Divider is at about x=6,003,700 EMU.
for shape in list(cur_slide.shapes):
    if shape.name not in keep_names and shape.left < 6000000:
        cur_slide.shapes._spTree.remove(shape._element)

copy_names = {
    '文本框 4', '文本框 5', '文本框 6', '文本框 9', '文本框 13', '矩形 14',
    'Rounded Rectangle 101', 'Rounded Rectangle 102', 'Rounded Rectangle 103'
}
for shape in src_slide.shapes:
    if shape.name in copy_names:
        cur_slide.shapes._spTree.append(deepcopy(shape._element))

cur.save(str(current))
print(current)
print(backup)
