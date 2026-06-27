from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

pptx_path = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹.pptx')
prs = Presentation(str(pptx_path))

# Delete unrelated slides, keep cover + one body page.
while len(prs.slides) > 2:
    r_id = prs.slides._sldIdLst[-1].rId
    prs.part.drop_rel(r_id)
    del prs.slides._sldIdLst[-1]

# Cover page text.
cover = prs.slides[0]
cover_texts = [sh for sh in cover.shapes if hasattr(sh, 'text') and sh.text.strip()]
for sh in cover_texts:
    txt = sh.text.strip()
    if '装备开发部月会' in txt:
        sh.text = '6月装备开发部月会'
    elif '月' in txt and '日' in txt:
        sh.text = '6月24日'
    elif '电气技术课' in txt:
        sh.text = '电气技术课 范宇丹'

slide = prs.slides[1]

# Remove old right-side HB content and image, keep title/owner/divider/left template.
remove_names = {
    '文本框 39', '组合 10', '图片 11', '文本框 16'
}
for sh in list(slide.shapes):
    if sh.name in remove_names:
        slide.shapes._spTree.remove(sh._element)

# Helper functions.
def set_text(shape_name, text, font_size=None, bold=None):
    for sh in slide.shapes:
        if sh.name == shape_name:
            sh.text = text
            for p in sh.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    if font_size is not None:
                        run.font.size = Pt(font_size)
                    if bold is not None:
                        run.font.bold = bold
            return sh
    raise RuntimeError(f'shape not found: {shape_name}')

def set_all_runs(shape, font_size=None, bold=None):
    if not hasattr(shape, 'text_frame'):
        return
    for p in shape.text_frame.paragraphs:
        for run in p.runs:
            if font_size is not None:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold

# Title and left body.
set_text('矩形 66', 'MINI优化', 24, True)
set_text('文本框 1', '责任人：范宇丹', 14, None)
set_text('文本框 4', '◆ 项目背景', 18, True)
set_text('文本框 5', 'MINI夹具闭合后回检功能在现场验证中发现：双臂/多工位回检目标存在顺序与缓存状态不一致风险；中间位置回检NG后暂停再继续时，剩余目标可能被跳过并直接进入后续测试。', 13, None)
set_text('文本框 6', '◆ 项目进展', 18, True)
progress = (
    '一：回检目标一致性优化：\n'
    '• 回检前按工站与模组位缓存目标，避免左右臂目标丢失\n'
    '• 区分 WS_ID + Num，防止不同工站同号模组互相覆盖\n'
    '• 未加载 WsModAfterClose_Shp 的相机侧按目标跳过并输出日志\n\n'
    '二：定时回检流程优化：\n'
    '• 回检轮次按目标逐个推进，避免工站1完成后关闭全局检测\n'
    '• 记录已完成目标，后续工站仍可继续进入本轮回检\n\n'
    '三：异常恢复优化：\n'
    '• 单点回检成功后仅移除当前目标\n'
    '• NG/暂停/异常打断时保留剩余目标，恢复后继续检查\n'
    '• 新增中断保留日志，便于现场复盘'
)
shape = set_text('文本框 9', progress, 11, None)
shape.text_frame.word_wrap = True
set_text('文本框 13', '效果：夹具合盖回检从“整批缓存一次性清空”调整为“逐点完成、异常保留”，降低漏检风险，提升多工位回检连续性和现场问题追溯效率。', 12, None)
set_text('矩形 14', '当前状态：已完成代码修复与日志级验证，未按项目规则默认编译；后续由现场版本编译后结合NG暂停场景复测。', 12, None)

# Update left process cards.
set_text('Rounded Rectangle 101', '目标缓存\nWS_ID + Num', 12, True)
set_text('Rounded Rectangle 102', '逐点执行\n完成即移除', 12, True)
set_text('Rounded Rectangle 103', '异常恢复\n保留未完成', 12, True)

# Keep right side blank intentionally. Add a tiny invisible note in speaker? Avoid adding visible text.
prs.save(str(pptx_path))
print(pptx_path)
