from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

pptx_path = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹.pptx')
prs = Presentation(str(pptx_path))
slide = prs.slides[1]

def find(name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise RuntimeError(name)

def write(shape, text, size=None, bold=None, color=None, align=PP_ALIGN.LEFT):
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = align
        for r in p.runs:
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color

write(find('矩形 66'), 'MINI优化', 24, True)
write(find('文本框 4'), '◆ 项目目的', 18, True)
write(find('文本框 5'), 'MINI一站式测试中，夹具闭合后需要通过上相机调用 WsModAfterClose_Shp，对合盖状态进行回检，确认通光孔/夹具位置无遮挡，避免夹具闭合异常流入后续测试。', 12)
write(find('文本框 6'), '◆ 问题与处理', 18, True)

progress = (
    '一：现场问题\n'
    '• 双臂/多工位运行时，回检目标按上下料缓存顺序执行，未按工位闭环管理\n'
    '• VPP只在单侧相机加载时，另一侧目标可能被跳过或无法单臂准确触发\n'
    '• 定时回检完成一个工站后，后续工站存在被全局开关提前结束的风险\n'
    '• 中间位置回检NG后暂停再继续，缓存被清空，剩余目标可能未回检即进入测试\n\n'
    '二：解决办法\n'
    '• 以 WS_ID + Num 作为回检目标唯一键，避免不同工站同号模组互相覆盖\n'
    '• 回检前确认目标相机是否加载 WsModAfterClose_Shp，未加载时按目标记录跳过日志\n'
    '• 定时回检改为按目标逐点推进，已完成目标单独记录，后续工站继续本轮检查\n'
    '• 单点成功后只移除当前目标；NG/暂停/异常打断时保留剩余目标，恢复后继续回检'
)
prog = find('文本框 9')
prog.left = Emu(180000)
prog.top = Emu(1810000)
prog.width = Emu(5450000)
prog.height = Emu(2700000)
write(prog, progress, 9.7)

result = find('文本框 13')
result.left = Emu(180000)
result.top = Emu(4580000)
result.width = Emu(5450000)
result.height = Emu(470000)
write(result, '结果：夹具闭合回检形成“目标确认、逐点执行、异常保留、恢复续检”的闭环，降低漏检和误跳过风险。', 11.5)

status = find('矩形 14')
status.left = Emu(180000)
status.top = Emu(5100000)
status.width = Emu(5450000)
status.height = Emu(440000)
write(status, '当前状态：代码修复与日志级验证已完成；待现场版本编译后，重点复测NG暂停恢复与多工位定时回检。', 11.5)

cards = [
    ('Rounded Rectangle 101', 228600, '目标确认\nWS_ID + Num'),
    ('Rounded Rectangle 102', 2077300, '逐点执行\n完成即移除'),
    ('Rounded Rectangle 103', 4040500, '异常续检\n保留剩余'),
]
for name, left, text in cards:
    sh = find(name)
    sh.left = Emu(left)
    sh.top = Emu(5580000)
    sh.width = Emu(1500000)
    sh.height = Emu(620000)
    write(sh, text, 12, True, RGBColor(31, 78, 121), PP_ALIGN.CENTER)

prs.save(str(pptx_path))
print(pptx_path)
