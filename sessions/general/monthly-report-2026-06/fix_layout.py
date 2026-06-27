from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

pptx_path = Path(r'E:\电气课室月报\2026年度月报\电气技术课6月月报-范宇丹.pptx')
prs = Presentation(str(pptx_path))
slide = prs.slides[1]

# Remove old mid-page process group from May template.
for sh in list(slide.shapes):
    if sh.name in {'组合 8', '组合 10', '图片 11', '文本框 16', '文本框 39'}:
        slide.shapes._spTree.remove(sh._element)

def find(name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise RuntimeError(name)

def write(shape, text, size=None, bold=None, color=None):
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.LEFT
        for r in p.runs:
            if size is not None:
                r.font.size = Pt(size)
            if bold is not None:
                r.font.bold = bold
            if color is not None:
                r.font.color.rgb = color

# Stable left layout.
write(find('矩形 66'), 'MINI优化', 24, True)
write(find('文本框 1'), '责任人：范宇丹', 14)
write(find('文本框 4'), '◆ 项目背景', 18, True)
write(find('文本框 5'), 'MINI夹具闭合后回检在双臂/多工位场景下，存在目标顺序、缓存状态与暂停恢复不一致风险；中间位置NG后恢复运行时，剩余目标可能被跳过。', 12)
write(find('文本框 6'), '◆ 项目进展', 18, True)

progress = (
    '一：目标一致性优化\n'
    '• 按工站与模组位缓存回检目标，区分 WS_ID + Num\n'
    '• 未加载 WsModAfterClose_Shp 的相机侧按目标跳过并记录日志\n\n'
    '二：定时回检流程优化\n'
    '• 回检轮次按目标逐点推进，避免单工站完成后关闭全局检测\n'
    '• 已完成目标单独记录，后续工站仍可继续本轮回检\n\n'
    '三：异常恢复优化\n'
    '• 单点成功后仅移除当前目标\n'
    '• NG/暂停/异常打断时保留剩余目标，恢复后继续检查'
)
prog = find('文本框 9')
prog.left = Emu(180000)
prog.top = Emu(1810000)
prog.width = Emu(5450000)
prog.height = Emu(2350000)
write(prog, progress, 10.2)

result = find('文本框 13')
result.left = Emu(180000)
result.top = Emu(4210000)
result.width = Emu(5450000)
result.height = Emu(470000)
write(result, '效果：夹具回检由“整批缓存一次性清空”调整为“逐点完成、异常保留”，降低漏检风险，提升多工位回检连续性。', 11.5)

status = find('矩形 14')
status.left = Emu(180000)
status.top = Emu(4770000)
status.width = Emu(5450000)
status.height = Emu(620000)
write(status, '当前状态：已完成代码修复与日志级验证；后续由现场版本编译后，结合NG暂停场景复测。', 11.5)

cards = [
    ('Rounded Rectangle 101', 228600, '目标缓存\nWS_ID + Num'),
    ('Rounded Rectangle 102', 2077300, '逐点执行\n完成即移除'),
    ('Rounded Rectangle 103', 4040500, '异常恢复\n保留未完成'),
]
for name, left, text in cards:
    sh = find(name)
    sh.left = Emu(left)
    sh.top = Emu(5580000)
    sh.width = Emu(1500000)
    sh.height = Emu(620000)
    write(sh, text, 12, True, RGBColor(31, 78, 121))
    for p in sh.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

prs.save(str(pptx_path))
print(pptx_path)
